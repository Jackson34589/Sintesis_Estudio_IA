import io
import base64
import hashlib
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
import fitz  # pymupdf
import anthropic
from fastapi import APIRouter, UploadFile, File, HTTPException
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

router = APIRouter()
client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

EXTRACT_PROMPT = """Extrae TODO el contenido visible exactamente como aparece.
Incluye:
- Todo el texto normal y párrafos
- Tablas completas (formato: columna1 | columna2 | columna3)
- Texto dentro de imágenes, figuras, gráficos, diagramas y esquemas
- Títulos, subtítulos, pies de figura, leyendas, etiquetas
- Números, unidades, fórmulas, referencias

No interpretes ni resumas — solo extrae el contenido tal cual aparece."""


def pix_to_b64(pix: fitz.Pixmap) -> str:
    return base64.standard_b64encode(pix.tobytes("png")).decode("utf-8")


def bytes_to_b64(img_bytes: bytes) -> str:
    return base64.standard_b64encode(img_bytes).decode("utf-8")


def extract_with_vision(b64: str) -> str:
    message = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}},
                {"type": "text", "text": EXTRACT_PROMPT},
            ],
        }],
    )
    return message.content[0].text.strip()


def render_page(page: fitz.Page, dpi: int) -> fitz.Pixmap:
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    return page.get_pixmap(matrix=mat)


MIN_IMAGE_DIM = 150  # píxeles mínimos por lado para considerar una imagen significativa


def extract_embedded_images_from_pdf(doc: fitz.Document) -> list:
    """Extrae imágenes embebidas individualmente del PDF (figuras, imágenes radiológicas, etc.).
    Deduplica por xref Y por hash del contenido para evitar repeticiones."""
    extracted = []
    seen_xrefs: set = set()
    seen_hashes: set = set()
    fig_count = 0

    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                pix = fitz.Pixmap(doc, xref)
                # Ignorar íconos y decoraciones pequeñas
                if pix.width < MIN_IMAGE_DIM or pix.height < MIN_IMAGE_DIM:
                    continue
                # Convertir CMYK u otros espacios de color no-RGB a RGB
                if pix.n - pix.alpha > 3:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                img_bytes = pix.tobytes("png")
                # Dedup por contenido: hash de los primeros 8 KB es suficiente
                content_hash = hashlib.md5(img_bytes[:8192]).hexdigest()
                if content_hash in seen_hashes:
                    continue
                seen_hashes.add(content_hash)
                b64 = base64.standard_b64encode(img_bytes).decode("utf-8")
                fig_count += 1
                extracted.append({
                    "data": b64,
                    "label": f"Figura {fig_count} (p.{page.number + 1})"
                })
            except Exception:
                pass

    return extracted


MAX_PDF_PAGES = 30   # cap to avoid memory spikes on huge PDFs
MAX_EMBEDDED_IMAGES = 8  # enough for a typical radiology document


def extract_pdf(content: bytes) -> dict:
    doc = fitz.open(stream=content, filetype="pdf")
    n_pages = min(len(doc), MAX_PDF_PAGES)

    # Submit pages to the thread pool one at a time so we never hold
    # all rendered bitmaps in memory simultaneously.
    page_texts: dict[int, str] = {}
    page_gallery: dict[int, str] = {}

    with ThreadPoolExecutor(max_workers=2) as pool:  # 2 workers = lower peak RAM
        futures: dict = {}
        for i in range(n_pages):
            page = doc[i]

            pix_v = render_page(page, dpi=100)   # 120→100: ~30 % smaller payload
            vision_b64 = pix_to_b64(pix_v)
            pix_v = None                          # release raw pixmap immediately

            pix_g = render_page(page, dpi=60)    # 72→60: gallery only needs low-res
            page_gallery[i] = pix_to_b64(pix_g)
            pix_g = None

            future = pool.submit(extract_with_vision, vision_b64)
            futures[future] = i
            vision_b64 = None                     # thread holds the only reference now

        for future in as_completed(futures):
            num = futures[future]
            try:
                page_texts[num] = future.result()
            except Exception:
                page_texts[num] = ""

    embedded_images = extract_embedded_images_from_pdf(doc)
    # Cap embedded images to avoid sending huge payloads to Claude
    embedded_images = embedded_images[:MAX_EMBEDDED_IMAGES]
    doc.close()

    pages_text = [page_texts[i] for i in range(n_pages) if page_texts.get(i)]
    page_images = [
        {"data": page_gallery[i], "label": f"Página {i + 1}"}
        for i in range(n_pages)
    ]

    if not pages_text:
        raise HTTPException(status_code=422, detail="No se pudo extraer contenido del PDF.")

    images = embedded_images if embedded_images else page_images

    return {
        "text": "\n\n".join(pages_text),
        "pages": len(pages_text),
        "images": images,
        "method": "vision",
    }


def extract_docx(content: bytes) -> dict:
    doc = Document(io.BytesIO(content))
    parts = []
    doc_images = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]
        if tag == "p":
            para = next((p for p in doc.paragraphs if p._element is element), None)
            if para and para.text.strip():
                parts.append(para.text.strip())
        elif tag == "tbl":
            tbl = next((t for t in doc.tables if t._element is element), None)
            if tbl:
                rows = [" | ".join(c.text.strip() for c in row.cells) for row in tbl.rows]
                parts.append("\n".join(rows))

    # Extract embedded images → Vision + gallery
    img_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_bytes = rel.target_part.blob
                img_doc = fitz.open(stream=img_bytes, filetype="image")
                pix = img_doc[0].get_pixmap()
                b64 = pix_to_b64(pix)
                img_doc.close()
                img_count += 1
                vision_text = extract_with_vision(b64)
                if vision_text:
                    parts.append(f"[Figura {img_count}]\n{vision_text}")
                doc_images.append({"data": b64, "label": f"Figura {img_count}"})
            except Exception:
                pass

    if not parts:
        raise HTTPException(status_code=422, detail="No se pudo extraer texto del archivo Word.")
    return {
        "text": "\n\n".join(parts),
        "pages": len(doc.paragraphs),
        "images": doc_images,
        "method": "docx+vision",
    }


def extract_pptx(content: bytes) -> dict:
    prs = Presentation(io.BytesIO(content))
    slides_text = []
    slide_images = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    slide_parts.append(text)

            if shape.has_table:
                rows = [" | ".join(c.text.strip() for c in row.cells) for row in shape.table.rows]
                slide_parts.append("\n".join(rows))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    img_doc = fitz.open(stream=img_bytes, filetype="image")
                    pix = img_doc[0].get_pixmap()
                    b64 = pix_to_b64(pix)
                    img_doc.close()
                    vision_text = extract_with_vision(b64)
                    if vision_text:
                        slide_parts.append(f"[Imagen]\n{vision_text}")
                    slide_images.append({"data": b64, "label": f"Diapositiva {i} — imagen"})
                except Exception:
                    pass

        if slide_parts:
            slides_text.append(f"[Diapositiva {i}]\n" + "\n\n".join(slide_parts))

    if not slides_text:
        raise HTTPException(status_code=422, detail="No se pudo extraer contenido del PowerPoint.")
    return {
        "text": "\n\n".join(slides_text),
        "pages": len(prs.slides),
        "images": slide_images,
        "method": "pptx+vision",
    }


@router.post("/extract-text")
async def extract_text(file: UploadFile = File(...)):
    filename = (file.filename or "").lower()
    content = await file.read()

    try:
        if filename.endswith(".txt"):
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                text = content.decode("latin-1")
            return {"text": text, "pages": None, "images": [], "method": "direct"}

        elif filename.endswith(".pdf"):
            return extract_pdf(content)

        elif filename.endswith(".docx"):
            return extract_docx(content)

        elif filename.endswith(".pptx"):
            return extract_pptx(content)

        elif filename.endswith(".doc"):
            raise HTTPException(status_code=415, detail="Formato .doc no soportado. Guarda como .docx e inténtalo de nuevo.")

        elif filename.endswith(".ppt"):
            raise HTTPException(status_code=415, detail="Formato .ppt no soportado. Guarda como .pptx e inténtalo de nuevo.")

        else:
            raise HTTPException(status_code=415, detail="Formato no soportado. Usa: .pdf, .docx, .pptx o .txt")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Error al procesar el archivo: {str(e)}")
