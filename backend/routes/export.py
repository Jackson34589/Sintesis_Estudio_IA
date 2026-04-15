import os
import re
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt

router = APIRouter()

EXPORTS_DIR = "exports"
os.makedirs(EXPORTS_DIR, exist_ok=True)


class ExportRequest(BaseModel):
    synthesis: str
    title: str = "Síntesis RadioSíntesis AI"


def parse_sections(text: str) -> list[dict]:
    """Split synthesis into title/content pairs for slides."""
    sections = []
    current_title = "Introducción"
    current_lines: list[str] = []

    for line in text.splitlines():
        stripped = line.strip()
        # Markdown headings or emoji-prefixed headings become slide titles
        if re.match(r"^#{1,3}\s+", stripped) or re.match(r"^[🎯📌🔬💡⚠️✅]\s*\*{0,2}.+", stripped):
            if current_lines:
                sections.append({"title": current_title, "content": "\n".join(current_lines).strip()})
                current_lines = []
            current_title = re.sub(r"^#{1,3}\s+", "", stripped).strip("*").strip()
        else:
            if stripped:
                current_lines.append(stripped)

    if current_lines:
        sections.append({"title": current_title, "content": "\n".join(current_lines).strip()})

    return sections or [{"title": "Síntesis", "content": text}]


@router.post("/export-pptx")
async def export_pptx(payload: ExportRequest):
    try:
        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = payload.title
        slide.placeholders[1].text = "Generado por RadioSíntesis AI"

        sections = parse_sections(payload.synthesis)
        content_layout = prs.slide_layouts[1]

        for section in sections:
            sl = prs.slides.add_slide(content_layout)
            sl.shapes.title.text = section["title"][:80]
            tf = sl.placeholders[1].text_frame
            tf.word_wrap = True
            # Limit content per slide to avoid overflow
            content = section["content"][:800]
            tf.text = content

        safe_title = re.sub(r"[^\w\s-]", "", payload.title).strip().replace(" ", "_")
        output_path = os.path.join(EXPORTS_DIR, f"{safe_title}.pptx")
        prs.save(output_path)
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{safe_title}.pptx",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al exportar PPTX: {str(e)}")
